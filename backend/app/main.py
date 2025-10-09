# app/main.py
import hashlib
import io
import os
import time
from contextlib import asynccontextmanager
from typing import List, Optional
from fastapi.responses import StreamingResponse
import openpyxl
from openpyxl.styles import Alignment, Font, PatternFill

from app import models, schemas
from app.analysis import analyze_xml  # ← app.schemas に正規化して返す実装に統一
from app.auth import get_current_user
# 既存DB CRUD（改修なし前提）
from app.crud import bulk_create_analysis_items, create_analysis, create_file
from app.crud import delete_file as delete_file_db
from app.crud import (get_analysis_with_items, get_file, get_latest_analysis,
                      list_analyses_by_file)
from app.crud import list_files as list_files_db
from app.db import Base, engine, get_db
from app.models import AnalysisItemRow
from app.pptx_parser import PptxConverter
from fastapi import (Depends, FastAPI, File, Form, HTTPException, Response,
                     UploadFile)
from fastapi.middleware.cors import CORSMiddleware
# Google / Vertex AI (Gemini)
from google import genai
from google.api_core.exceptions import NotFound
# Google Cloud
from google.cloud import secretmanager, storage
from pptx import Presentation
from pydantic import BaseModel
from sqlalchemy import text
from sqlalchemy.orm import Session


# =========================
# Secret Manager
# =========================
def access_secret_version(project_id: str, secret_id: str, version_id: str = "latest") -> str:
    client = secretmanager.SecretManagerServiceClient()
    name = f"projects/{project_id}/secrets/{secret_id}/versions/{version_id}"
    response = client.access_secret_version(request={"name": name})
    return response.payload.data.decode("UTF-8")


# =========================
# GCS Client
# =========================
GCS_BUCKET_NAME = os.getenv("GCS_BUCKET_NAME")
if not GCS_BUCKET_NAME:
    raise RuntimeError("GCS_BUCKET_NAME environment variable not set.")
storage_client = storage.Client()
bucket = storage_client.bucket(GCS_BUCKET_NAME)


# =========================
# Helpers
# =========================
def _sha256_bytes(b: bytes) -> str:
    h = hashlib.sha256()
    h.update(b)
    return h.hexdigest()


def extract_pptx_text_from_file(file_obj) -> str:
    """in-memory の PPTX から素朴にテキスト抽出（デバッグ用）"""
    presentation = Presentation(file_obj)
    full_text = []
    for i, slide in enumerate(presentation.slides):
        slide_text = [f"[スライド {i + 1}]"]
        has_text = False
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text and run.text.strip():
                        slide_text.append(f"  - {run.text.strip()}")
                        has_text = True
        if not has_text:
            slide_text.append("  (テキストなし)")
        full_text.append("\n".join(slide_text))
    return "\n\n".join(full_text)


async def extract_and_print_pptx_text(file: UploadFile):
    """アップロードされた PPTX の中身をコンソールに表示（ファイルポインタは戻す）"""
    if not file.filename or not file.filename.lower().endswith(".pptx"):
        return
    print(f"\n--- 📄 PPTX '{file.filename}' テキスト抽出 ---")
    try:
        presentation = Presentation(file.file)
        for i, slide in enumerate(presentation.slides):
            print(f"\n[スライド {i + 1}]")
            has_text = False
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            print(f"  - {run.text.strip()}")
                            has_text = True
            if not has_text:
                print("  (テキストなし)")
    except Exception as e:
        print(f"❌ PPTX解析エラー: {e}")
    finally:
        await file.seek(0)
        print("--- ✅ 抽出終了 ---\n")


# =========================
# Lifespan
# =========================
@asynccontextmanager
async def lifespan(app: FastAPI):
    # DB起動待機 & テーブル作成
    max_tries = 30
    for i in range(max_tries):
        try:
            with engine.connect() as conn:
                conn.execute(text("SELECT 1"))
            break
        except Exception:
            print(f"⏳ Waiting for database... ({i + 1}/{max_tries})")
            time.sleep(1)
    else:
        raise RuntimeError("❌ Database connection failed after waiting.")

    print("✅ Database connected. Creating tables if not exist...")
    Base.metadata.create_all(bind=engine)
    yield


# =========================
# FastAPI app & CORS
# =========================
app = FastAPI(lifespan=lifespan)
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        # あなたの公開サイトからのアクセスを許可
        "https://storage.googleapis.com/minapp-frontend-bucket-474213", 
        
        # あなたのパソコンでの開発・テストからのアクセスを許可
        "http://localhost:5173",
        "http://127.0.0.1:5173",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# =========================
# Me / Health
# =========================
class MeResp(BaseModel):
    uid: str
    email: str | None = None
    emailVerified: bool = False


@app.get("/api/me", response_model=MeResp)
def get_me(response: Response, user=Depends(get_current_user)):
    response.headers["Cache-Control"] = "no-store"
    return {"uid": user["uid"], "email": user.get("email"), "emailVerified": user.get("verified", False)}


@app.get("/api/")
def read_root():
    return {"message": "Hello from MinApp Backend with DB & CORS!"}


@app.get("/api/health")
def health():
    return {"status": "ok"}


@app.get("/api/health/deep")
def health_deep(db: Session = Depends(get_db)):
    db_ok = True
    try:
        db.execute(text("SELECT 1"))
    except Exception:
        db_ok = False
    gcs_ok = True
    try:
        _ = list(bucket.list_blobs(prefix="__healthcheck__/"))
    except Exception:
        gcs_ok = False
    status = "ok" if (db_ok and gcs_ok) else "degraded"
    return {"status": status, "db": db_ok, "gcs": gcs_ok}


# =========================
# Files: GCS list (併存)
# =========================
@app.get("/api/files", response_model=List[schemas.GcsFile])
def list_files(user=Depends(get_current_user)):
    prefix = f"{user['uid']}/"
    blobs = bucket.list_blobs(prefix=prefix)
    return [
        {"name": b.name.replace(prefix, ""), "updated": b.updated, "size": b.size, "path": b.name}
        for b in blobs
        if b.name != prefix
    ]


# =========================
# Files: Upload (GCS保存 + DB登録) ← 旧APIの動線を引き継ぐ
# =========================
@app.post("/api/files/upload")
async def upload_file(user=Depends(get_current_user), file: UploadFile = File(...), db: Session = Depends(get_db)):
    if not file or not (file.filename or "").lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="pptxファイルのみ対応しています")

    data = await file.read()  # DB記録のため先に読み切る
    await file.seek(0)
    destination = f"{user['uid']}/{file.filename}"
    blob = bucket.blob(destination)
    try:
        # デバッグ出力（任意）
        await extract_and_print_pptx_text(file)
        # GCSへ
        blob.upload_from_file(file.file, content_type=file.content_type)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"アップロード失敗: {e}")

    digest = _sha256_bytes(data)
    rec = create_file(
        db,
        user_id=user["uid"],
        filename=file.filename,
        path=destination,        # GCSのパスをDBへ
        sha256=digest,
        size_bytes=len(data),
    )
    return {
        "file_id": rec.id,
        "filename": rec.filename,
        "size_bytes": rec.size_bytes,
        "sha256": rec.sha256,
        "path": destination,
    }


# =========================
# Files: DBベース一覧（旧互換の形で返す）
# =========================
@app.get("/api/files-db")
def get_files_db(db: Session = Depends(get_db), user=Depends(get_current_user)):
    files = list_files_db(db, user_id=user["uid"])
    result = []
    for f in files:
        analyses = list_analyses_by_file(db, file_id=f.id, user_id=user["uid"])
        status = "success" if len(analyses) > 0 else "pending"
        result.append(
            {
                "id": str(f.id),
                "file": None,
                "name": f.filename,
                "size": f.size_bytes,
                "uploadDate": f.created_at.isoformat(),
                "status": status,
                "analysisResult": [],
                "error": None,
                "isBasisAugmented": False,
                "augmentationStatus": "idle",
            }
        )
    return result


# =========================
# Files: GCSファイル読み取り（簡易）
# =========================
@app.post("/api/files/read", response_model=schemas.FileContent)
def read_file_content(file_path: schemas.GcsFilePath, user=Depends(get_current_user)):
    if not file_path.path.startswith(f"{user['uid']}/"):
        raise HTTPException(status_code=403, detail="Access denied")
    blob = bucket.blob(file_path.path)
    if not blob.exists():
        raise HTTPException(status_code=404, detail="ファイルが見つかりません")
    try:
        file_bytes = blob.download_as_bytes()
        if file_path.path.lower().endswith(".pptx"):
            content = extract_pptx_text_from_file(io.BytesIO(file_bytes))
        else:
            content = "このファイル形式は現在、内容の読み取りに対応していません。"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"ファイルの読み取り中にエラーが発生しました: {e}")
    return {"content": content}


# =========================
# Files: DB基点の削除（GCS→DB）
# =========================
@app.delete("/api/files/{file_id}")
def remove_file_db(file_id: int, db: Session = Depends(get_db), user=Depends(get_current_user)):
    f = get_file(db, file_id)
    if not f or f.user_id != user["uid"]:
        raise HTTPException(404, "not found")
    # GCS削除（無くても続行）
    try:
        bucket.blob(f.path).delete()
    except Exception:
        pass
    delete_file_db(db, file_id)
    return {"ok": True}


# =========================
# PPTX → XML（Upload または GCS パス）
# =========================
@app.post("/api/pptx/xml")
async def pptx_to_xml(
    file: UploadFile = File(None),
    gcs_path: Optional[str] = Form(None),  # 例: "<uid>/slides/deck.pptx"
    pretty: Optional[bool] = Form(True),
    user=Depends(get_current_user),
):
    if not file and not gcs_path:
        raise HTTPException(400, "file または gcs_path のいずれかを指定してください")

    if file:
        if not (file.filename or "").lower().endswith(".pptx"):
            raise HTTPException(400, "pptxファイルのみ対応しています")
        xml_str = PptxConverter.convert_to_xml(file, pretty=bool(pretty))
        return Response(content=xml_str, media_type="application/xml")

    if not gcs_path.startswith(f"{user['uid']}/"):
        raise HTTPException(403, "Access denied")
    if not gcs_path.lower().endswith(".pptx"):
        raise HTTPException(400, "pptxファイルのみ対応しています")

    blob = bucket.blob(gcs_path)
    if not blob.exists():
        raise HTTPException(status_code=404, detail="ファイルが見つかりません")

    data = blob.download_as_bytes()
    xml_str = PptxConverter.convert_to_xml(io.BytesIO(data), pretty=bool(pretty))
    return Response(content=xml_str, media_type="application/xml")


# =========================
# 解析（file_id または gcs_path）: 実行して DB 保存
# =========================
def _mock_items_from_xml(_: str) -> list[schemas.AnalysisItem]:
    return [
        schemas.AnalysisItem(
            slideNumber=1,
            category="表現",
            basis="1",
            issue="サンプル：表現が断定的です。",
            suggestion="サンプル：より中立な言い回しに調整してください。",
            correctionType="任意",
        )
    ]


@app.post("/api/analyze", response_model=List[schemas.AnalysisItem])
async def analyze_endpoint(
    file_id: Optional[int] = Form(None),
    gcs_path: Optional[str] = Form(None),     # どちらか必須
    rules: Optional[str] = Form(None),
    mode: Optional[str] = Form(None),         # auto | mock | llm（省略時: auto）
    db: Session = Depends(get_db),
    user=Depends(get_current_user),
):
    if not file_id and not gcs_path:
        raise HTTPException(400, "file_id または gcs_path のいずれかを指定してください")

    # file_id が来たら DB から GCS パスを引く
    if file_id:
        f = get_file(db, file_id)
        if not f or f.user_id != user["uid"]:
            raise HTTPException(404, "file not found")
        gcs_path = f.path

    assert gcs_path is not None
    if not gcs_path.startswith(f"{user['uid']}/"):
        raise HTTPException(403, "Access denied")
    if not gcs_path.lower().endswith(".pptx"):
        raise HTTPException(400, "pptxファイルのみ対応しています")

    # 1) GCS → XML
    blob = bucket.blob(gcs_path)
    if not blob.exists():
        raise HTTPException(status_code=404, detail="ファイルが見つかりません")
    try:
        data = blob.download_as_bytes()
        xml_str = PptxConverter.convert_to_xml(io.BytesIO(data), pretty=False)
    except Exception as e:
        raise HTTPException(500, f"pptx→xml 変換中にエラー: {e}")

    # 2) モード判定
    req_mode = (mode or "auto").lower()
    if req_mode not in {"auto", "mock", "llm"}:
        raise HTTPException(400, "mode は auto/mock/llm のいずれかを指定してください")

    # 3) 実行
    items: list[schemas.AnalysisItem]
    model_name = "gemini-2.5-flash"
    if req_mode == "mock":
        items = _mock_items_from_xml(xml_str)
        model_name = "mock"
    else:
        # Secret Manager から API キー取得（常に渡して安定化）
        project_id = os.getenv("GOOGLE_PROJECT_ID")
        if not project_id:
            raise HTTPException(status_code=500, detail="GOOGLE_PROJECT_ID not set.")
        try:
            api_key = access_secret_version(project_id=project_id, secret_id="API_KEY")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Secret fetch failed: {e}")

        try:
            items = analyze_xml(xml_str, api_key=api_key)  # rules を使うなら analyze_xml 側で拡張
        except Exception:
            if req_mode == "llm":
                raise HTTPException(502, "LLM 解析に失敗しました")
            items = _mock_items_from_xml(xml_str)
            model_name = "mock"

    # 4) 補正
    for i in items:
        if getattr(i, "correctionType", None) is None:
            i.correctionType = "任意"

    # 5) DB 保存（file_id がある場合のみ）
    if file_id:
        payload = [i.model_dump() for i in items]
        analysis = create_analysis(
            db,
            user_id=user["uid"],
            file_id=file_id,
            model=model_name,
            rules_version=None,   # rules をバージョン管理するならここで設定
            result_json=payload,
        )
        bulk_create_analysis_items(db, analysis_id=analysis.id, items=items)

    return items


# =========================
# 解析の一覧/詳細/最新/アイテムCRUD（旧API互換）
# =========================
@app.get("/api/files/{file_id}/analyses")
def list_analyses_api(file_id: int, db: Session = Depends(get_db), user=Depends(get_current_user)):
    f = get_file(db, file_id)
    if not f or f.user_id != user["uid"]:
        raise HTTPException(404, "not found")
    lst = list_analyses_by_file(db, file_id=file_id, user_id=user["uid"])
    return [
        {"id": a.id, "created_at": str(a.created_at), "model": a.model, "status": a.status, "items_count": len(a.items)}
        for a in lst
    ]


@app.get("/api/analyses/{analysis_id}")
def get_analysis_api(analysis_id: int, db: Session = Depends(get_db), user=Depends(get_current_user)):
    a, rows = get_analysis_with_items(db, analysis_id)
    if not a or a.user_id != user["uid"]:
        raise HTTPException(404, "analysis not found")
    return {
        "id": a.id,
        "file_id": a.file_id,
        "created_at": str(a.created_at),
        "model": a.model,
        "status": a.status,
        "rules_version": a.rules_version,
        "result_json": a.result_json,
        "items": [
            {
                "id": r.id,
                "slideNumber": r.slide_number,
                "category": r.category,
                "basis": r.basis,
                "issue": r.issue,
                "suggestion": r.suggestion,
                "correctionType": r.correction_type,
            }
            for r in rows
        ],
    }


@app.get("/api/files/{file_id}/analyses/latest")
def get_latest_analysis_for_file_api(file_id: int, db: Session = Depends(get_db), user=Depends(get_current_user)):
    latest, items = get_latest_analysis(db, file_id=file_id, user_id=user["uid"])
    if not latest:
        raise HTTPException(404, "no analysis found")
    return {
        "id": latest.id,
        "created_at": str(latest.created_at),
        "model": latest.model,
        "status": latest.status,
        "items": [
            {
                "id": r.id,
                "slideNumber": r.slide_number,
                "category": r.category,
                "basis": r.basis,
                "issue": r.issue,
                "suggestion": r.suggestion,
                "correctionType": r.correction_type,
            }
            for r in items
        ],
    }


@app.post("/api/files/{file_id}/analyses/latest/items")
def add_item_to_latest_analysis_api(
    file_id: int,
    item: schemas.AnalysisItemCreate,
    db: Session = Depends(get_db),
    user=Depends(get_current_user),
):
    f = get_file(db, file_id)
    if not f or f.user_id != user["uid"]:
        raise HTTPException(404, "file not found")

    analyses = list_analyses_by_file(db, file_id=file_id, user_id=user["uid"])
    if analyses:
        a = analyses[0]
    else:
        a = create_analysis(
            db,
            user_id=user["uid"],
            file_id=file_id,
            model="manual-edit",
            rules_version=None,
            result_json=[],
        )

    row = AnalysisItemRow(
        analysis_id=a.id,
        slide_number=item.slideNumber,
        category=item.category,
        basis=item.basis,
        issue=item.issue,
        suggestion=item.suggestion,
        correction_type=item.correctionType or "任意",
    )
    db.add(row)
    db.commit()
    db.refresh(row)
    return {
        "id": row.id,
        "slideNumber": row.slide_number,
        "category": row.category,
        "basis": row.basis,
        "issue": row.issue,
        "suggestion": row.suggestion,
        "correctionType": row.correction_type,
    }


@app.patch("/api/analysis-items/{item_id}")
def update_analysis_item_api(
    item_id: int,
    patch: schemas.AnalysisItemUpdate,
    db: Session = Depends(get_db),
    user=Depends(get_current_user),
):
    row = db.get(AnalysisItemRow, item_id)
    if not row:
        raise HTTPException(404, "item not found")
    if not row.analysis or row.analysis.user_id != user["uid"]:
        raise HTTPException(403, "forbidden")

    if patch.slideNumber is not None:
        row.slide_number = patch.slideNumber
    if patch.category is not None:
        row.category = patch.category
    if patch.basis is not None:
        row.basis = patch.basis
    if patch.issue is not None:
        row.issue = patch.issue
    if patch.suggestion is not None:
        row.suggestion = patch.suggestion
    if patch.correctionType is not None:
        row.correction_type = patch.correctionType

    db.add(row)
    db.commit()
    db.refresh(row)
    return {
        "id": row.id,
        "slideNumber": row.slide_number,
        "category": row.category,
        "basis": row.basis,
        "issue": row.issue,
        "suggestion": row.suggestion,
        "correctionType": row.correction_type,
    }


@app.delete("/api/analysis-items/{item_id}")
def delete_analysis_item_api(
    item_id: int,
    db: Session = Depends(get_db),
    user=Depends(get_current_user),
):
    row = db.get(AnalysisItemRow, item_id)
    if not row:
        raise HTTPException(404, "item not found")
    if not row.analysis or row.analysis.user_id != user["uid"]:
        raise HTTPException(403, "forbidden")
    db.delete(row)
    db.commit()
    return {"ok": True}

# =========================
# Excel Export (追加)
# =========================
@app.get("/api/export/analysis/{analysis_id}/excel")
def export_analysis_to_excel(
    analysis_id: int,
    db: Session = Depends(get_db),
    user=Depends(get_current_user)
):
    analysis, items = get_analysis_with_items(db, analysis_id)
    if not analysis or analysis.user_id != user["uid"]:
        raise HTTPException(status_code=404, detail="Analysis not found")

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "AnalysisResults"

    # ヘッダー行を書き込み
    headers = ["スライド番号", "カテゴリ", "根拠", "指摘事項", "改善案", "修正の種類"]
    sheet.append(headers)

    # ヘッダーのスタイリング
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="4F81BD", end_color="4F81BD", fill_type="solid")
    for cell in sheet[1]:
        cell.font = header_font
        cell.fill = header_fill

    # データ行を書き込み & スタイリング
    wrap_alignment = Alignment(wrap_text=True, vertical='top')
    required_fill = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")

    for item in items:
        row_data = [
            item.slide_number,
            item.category,
            item.basis,
            item.issue,
            item.suggestion,
            item.correction_type
        ]
        sheet.append(row_data)
        
        # 「必須」の行に色を付ける
        if item.correction_type == "必須":
            for cell in sheet[sheet.max_row]:
                cell.fill = required_fill
    
    # 全てのデータセルにテキスト折り返しを設定
    for row in sheet.iter_rows(min_row=2):
        for cell in row:
            cell.alignment = wrap_alignment

    # 列幅の調整
    sheet.column_dimensions['A'].width = 12  # スライド番号
    sheet.column_dimensions['B'].width = 15  # カテゴリ
    sheet.column_dimensions['C'].width = 40  # 根拠
    sheet.column_dimensions['D'].width = 60  # 指摘事項
    sheet.column_dimensions['E'].width = 60  # 改善案
    sheet.column_dimensions['F'].width = 12  # 修正の種類

    excel_stream = io.BytesIO()
    workbook.save(excel_stream)
    excel_stream.seek(0)

    file_name = f"analysis_{analysis_id}.xlsx"
    headers = {
        'Content-Disposition': f'attachment; filename="{file_name}"'
    }
    return StreamingResponse(
        excel_stream,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers=headers
    )
